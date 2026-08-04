"""
Renders each slide to a PNG using Pillow so the layout can be eyeballed
without PowerPoint/LibreOffice. Approximate but faithful enough to catch
spacing, colour and collision problems.
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from PIL import Image, ImageDraw, ImageFont

EMU_IN = 914400
SCALE = 110  # px per inch
HERE = os.path.dirname(os.path.abspath(__file__))
PATH = os.path.join(HERE, "Alur-Kerja-WJW.pptx")
OUT = os.path.join(HERE, "preview")
os.makedirs(OUT, exist_ok=True)

FONT_DIRS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
]
REG = FONT_DIRS[0] if os.path.exists(FONT_DIRS[0]) else None
BOLD = FONT_DIRS[1] if os.path.exists(FONT_DIRS[1]) else REG


def font(pt, bold=False):
    px = max(7, int(pt * SCALE / 72))
    try:
        return ImageFont.truetype(BOLD if bold else REG, px)
    except Exception:
        return ImageFont.load_default()


def rgb(c):
    if c is None:
        return None
    try:
        return (c[0], c[1], c[2])
    except Exception:
        return None


def shape_fill(sp):
    try:
        if sp.fill.type is not None and sp.fill.type == 1:
            return rgb(sp.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def shape_line(sp):
    try:
        if sp.line.fill.type == 1:
            return rgb(sp.line.color.rgb)
    except Exception:
        pass
    return None


prs = Presentation(PATH)
W = int(prs.slide_width / EMU_IN * SCALE)
H = int(prs.slide_height / EMU_IN * SCALE)

for idx, slide in enumerate(prs.slides, start=1):
    img = Image.new("RGB", (W, H), (13, 17, 23))
    d = ImageDraw.Draw(img)

    for sp in slide.shapes:
        if sp.left is None:
            continue
        x = sp.left / EMU_IN * SCALE
        y = sp.top / EMU_IN * SCALE
        w = sp.width / EMU_IN * SCALE
        h = sp.height / EMU_IN * SCALE
        box = [x, y, x + w, y + h]

        fill = shape_fill(sp)
        line = shape_line(sp)
        st = sp.shape_type

        if fill or line:
            try:
                if sp.auto_shape_type == MSO_SHAPE.OVAL:
                    d.ellipse(box, fill=fill, outline=line, width=2)
                elif sp.auto_shape_type == MSO_SHAPE.DIAMOND:
                    cx, cy = x + w / 2, y + h / 2
                    d.polygon(
                        [(cx, y), (x + w, cy), (cx, y + h), (x, cy)],
                        fill=fill, outline=line,
                    )
                elif sp.auto_shape_type == MSO_SHAPE.RIGHT_ARROW:
                    mid = y + h / 2
                    th = h * 0.44
                    hd = min(w * 0.45, h * 0.9)
                    d.polygon(
                        [
                            (x, mid - th / 2), (x + w - hd, mid - th / 2),
                            (x + w - hd, y), (x + w, mid),
                            (x + w - hd, y + h), (x + w - hd, mid + th / 2),
                            (x, mid + th / 2),
                        ],
                        fill=fill,
                    )
                elif sp.auto_shape_type == MSO_SHAPE.DOWN_ARROW:
                    midx = x + w / 2
                    tw = w * 0.42
                    hd = min(h * 0.5, w * 1.2)
                    d.polygon(
                        [
                            (midx - tw / 2, y), (midx + tw / 2, y),
                            (midx + tw / 2, y + h - hd), (x + w, y + h - hd),
                            (midx, y + h), (x, y + h - hd),
                            (midx - tw / 2, y + h - hd),
                        ],
                        fill=fill,
                    )
                else:
                    r = min(14, h / 2, w / 2)
                    d.rounded_rectangle(box, radius=r, fill=fill,
                                        outline=line, width=2)
            except Exception:
                d.rectangle(box, fill=fill, outline=line)

        if not sp.has_text_frame:
            continue

        tf = sp.text_frame
        paras = [p for p in tf.paragraphs]
        total = ""
        for p in paras:
            total += "".join(r.text for r in p.runs) + "\n"
        if not total.strip():
            continue

        # vertical anchor
        line_hs = []
        for p in paras:
            pt = max([r.font.size.pt for r in p.runs if r.font.size] or [12])
            line_hs.append(pt * 1.35 * SCALE / 72)
        blk = sum(line_hs)
        anchor = tf.vertical_anchor
        if anchor == 3:  # MIDDLE
            cy = y + (h - blk) / 2
        elif anchor == 4:  # BOTTOM
            cy = y + h - blk
        else:
            cy = y + 3

        for p, lh in zip(paras, line_hs):
            txt = "".join(r.text for r in p.runs)
            if not txt:
                cy += lh
                continue
            run = p.runs[0]
            pt = run.font.size.pt if run.font.size else 12
            bold = bool(run.font.bold)
            col = rgb(run.font.color.rgb) if run.font.color and run.font.color.type is not None else (238, 242, 248)
            f = font(pt, bold)
            try:
                tw = d.textlength(txt, font=f)
            except Exception:
                tw = len(txt) * pt * 0.5 * SCALE / 72
            al = p.alignment
            if al == 2:  # CENTER
                tx = x + (w - tw) / 2
            elif al == 3:  # RIGHT
                tx = x + w - tw
            else:
                tx = x + 4
            d.text((tx, cy), txt, font=f, fill=col or (238, 242, 248))
            cy += lh

    img.save(os.path.join(OUT, f"slide{idx}.png"))
    print("rendered slide", idx)

print("done ->", OUT)
