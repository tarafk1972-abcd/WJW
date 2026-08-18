"""
Generates 'Alur-Kerja-WJW.pptx' — workflow diagrams for Warga Jaga Warga.

Run:  python3 docs/build_flow_ppt.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------- palette
BG = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x1A, 0x21, 0x30)
SURFACE2 = RGBColor(0x22, 0x2B, 0x3C)
LINE = RGBColor(0x2C, 0x36, 0x48)
TEXT = RGBColor(0xEE, 0xF2, 0xF8)
TEXT2 = RGBColor(0x9A, 0xA7, 0xBD)
TEXT3 = RGBColor(0x6B, 0x78, 0x91)
BRAND = RGBColor(0x2E, 0xC2, 0x7E)
DANGER = RGBColor(0xFF, 0x4D, 0x5E)
WARN = RGBColor(0xFF, 0xB5, 0x45)
INFO = RGBColor(0x58, 0xA6, 0xFF)
PURPLE = RGBColor(0xA3, 0x71, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

SW, SH = Emu(12192000), Emu(6858000)  # 16:9


def inches(v):
    return Emu(int(v * 914400))


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(s, x, y, w, h, text, size=14, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
            spacing=1.0):
    tb = s.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def box(s, x, y, w, h, title, body="", fill=SURFACE, border=LINE,
        title_color=TEXT, body_color=TEXT2, title_size=13, body_size=10,
        radius=0.10, shape=MSO_SHAPE.ROUNDED_RECTANGLE, border_w=1.25):
    sh = s.shapes.add_shape(shape, inches(x), inches(y), inches(w), inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = inches(0.10)
    tf.margin_top = tf.margin_bottom = inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    if body:
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.line_spacing = 0.95
            r2 = p2.add_run()
            r2.text = line
            r2.font.name = FONT
            r2.font.size = Pt(body_size)
            r2.font.color.rgb = body_color
    return sh


def arrow(s, x, y, w, color=TEXT3, h=0.052, head=0.42):
    """Horizontal arrow pointing right."""
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(x), inches(y),
                           inches(w), inches(h * 4))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.55
        a.adjustments[1] = head
    except (IndexError, KeyError):
        pass
    return a


def down_arrow(s, x, y, h, color=TEXT3, w=0.20):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, inches(x), inches(y),
                           inches(w), inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.45
        a.adjustments[1] = 0.62
    except (IndexError, KeyError):
        pass
    return a


def pill(s, x, y, w, h, text, fill, color=WHITE, size=9.5, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y),
                            inches(w), inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.5
    except (IndexError, KeyError):
        pass
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def header(s, kicker, title, accent=BRAND):
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(0.55),
                             inches(0.42), inches(0.075), inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    textbox(s, 0.78, 0.40, 9.5, 0.26, kicker, size=10.5, color=accent, bold=True)
    textbox(s, 0.78, 0.62, 11.0, 0.45, title, size=23, color=TEXT, bold=True)


def footnote(s, text, color=TEXT3, y=6.82):
    textbox(s, 0.58, y, 12.2, 0.3, text, size=9.5, color=color, italic=True)


def legend(s, items, x=0.58, y=6.62):
    """items: list of (label, color)"""
    cx = x
    for label, col in items:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(cx), inches(y + 0.045),
                               inches(0.11), inches(0.11))
        d.fill.solid()
        d.fill.fore_color.rgb = col
        d.line.fill.background()
        d.shadow.inherit = False
        textbox(s, cx + 0.17, y, 2.4, 0.22, label, size=9.5, color=TEXT2)
        cx += 0.20 + 0.10 + (len(label) * 0.062) + 0.34


# ================================================================ SLIDE 1
s = slide()
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(-1.6), inches(4.9),
                          inches(7.2), inches(4.2))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x11, 0x20, 0x1D)
glow.line.fill.background()
glow.shadow.inherit = False
glow.name = 'decor-bleed'

logo = box(s, 0.9, 1.66, 0.92, 0.92, "WJW", fill=BRAND,
           border=BRAND, title_color=RGBColor(0x04, 0x24, 0x17), title_size=20,
           radius=0.22)

textbox(s, 2.10, 1.62, 6.6, 0.75, "Warga Jaga Warga", size=38, color=TEXT, bold=True)
textbox(s, 2.13, 2.50, 6.6, 0.32, "Aplikasi keamanan lingkungan warga",
        size=15, color=BRAND, bold=True)

textbox(s, 0.92, 3.35, 7.6, 1.0,
        "Diagram Alur Kerja Aplikasi\n"
        "Dari pendaftaran warga sampai penanganan darurat.",
        size=13, color=TEXT2, spacing=1.35)

chips = [
    ("3 Cara Bergabung", INFO),
    ("Tombol Panik", DANGER),
    ("Persetujuan Admin", PURPLE),
    ("Jaringan Bantuan", BRAND),
]
cx = 0.92
for label, col in chips:
    w = 0.30 + len(label) * 0.088
    pill(s, cx, 4.60, w, 0.34, label, fill=SURFACE2, color=col, size=10.5)
    cx += w + 0.16

textbox(s, 0.92, 5.55, 8, 0.28,
        "Versi 1 (MVP)  ·  Bahasa Indonesia / English / Basa Sunda",
        size=11, color=TEXT3)

# side panel
box(s, 9.05, 1.75, 3.35, 3.9, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 9.32, 2.00, 2.9, 0.3, "ISI DECK", size=10, color=TEXT3, bold=True)
toc = [
    ("1", "Peta alur keseluruhan"),
    ("2", "Cara gabung / buat lingkungan"),
    ("3", "Persetujuan admin & peran"),
    ("4", "Alur tombol panik (MVP)"),
    ("5", "Penerima & respons"),
    ("6", "Peran & hak akses"),
]
ty = 2.42
for num, label in toc:
    pill(s, 9.32, ty - 0.015, 0.26, 0.26, num, fill=SURFACE2, color=BRAND, size=9.5)
    textbox(s, 9.70, ty + 0.005, 2.6, 0.26, label, size=10.5, color=TEXT2)
    ty += 0.52

footnote(s, "Dokumen ini menggambarkan aplikasi yang sudah berjalan, bukan rencana.", y=6.35)

# ================================================================ SLIDE 2
s = slide()
header(s, "01  ·  GAMBARAN UMUM", "Peta Alur Keseluruhan")

lanes = [
    ("MASUK", 0.58, INFO),
    ("VERIFIKASI", 3.62, PURPLE),
    ("PENGGUNAAN HARIAN", 6.66, BRAND),
    ("DARURAT", 9.70, DANGER),
]
for label, x, col in lanes:
    pill(s, x, 1.42, 2.72, 0.34, label, fill=SURFACE2, color=col, size=10.5)

# lane 1
box(s, 0.58, 2.00, 2.72, 0.72, "Buka aplikasi",
    "Pilih bahasa\n(ID / EN / SU)", fill=SURFACE, border=INFO)
down_arrow(s, 1.84, 2.80, 0.32, color=TEXT3)
box(s, 0.58, 3.20, 2.72, 1.00, "Gabung / Buat",
    "1. Kode undangan / QR\n2. Cari lingkungan\n3. Buat baru → Admin",
    fill=SURFACE, border=INFO)

# lane 2
box(s, 3.62, 2.00, 2.72, 0.72, "Antre persetujuan",
    "Status: menunggu", fill=SURFACE, border=PURPLE)
down_arrow(s, 4.88, 2.80, 0.32, color=TEXT3)
box(s, 3.62, 3.20, 2.72, 1.00, "Admin memutuskan",
    "Terima → Warga /\nSatpam / Admin\natau Tolak + alasan",
    fill=SURFACE, border=PURPLE)

# lane 3
box(s, 6.66, 2.00, 2.72, 0.72, "Sapaan pribadi",
    "“Apa kabar hari ini,\n<nama>?”", fill=SURFACE, border=BRAND)
down_arrow(s, 7.92, 2.80, 0.32, color=TEXT3)
box(s, 6.66, 3.20, 2.72, 1.00, "Fitur harian",
    "Lapor · Peta area · Tamu\nPatroli · Pengumuman\nJaringan bantuan",
    fill=SURFACE, border=BRAND)

# lane 4
box(s, 9.70, 2.00, 2.72, 0.72, "Tombol DARURAT",
    "Tahan 2 detik", fill=SURFACE, border=DANGER)
down_arrow(s, 10.96, 2.80, 0.32, color=TEXT3)
box(s, 9.70, 3.20, 2.72, 1.00, "Kirim peringatan",
    "GPS · profil · suara 15s\nfoto/video · waktu\n→ jaringan bantuan",
    fill=SURFACE, border=DANGER)

for x in (3.34, 6.38, 9.42):
    arrow(s, x, 2.19, 0.24, color=TEXT3)

# bottom strip
box(s, 0.58, 4.62, 11.84, 1.30, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 0.85, 4.82, 4.0, 0.26, "BERJALAN DI LATAR", size=10, color=TEXT3, bold=True)
bl = [
    ("Masa percobaan 14 hari", WARN),
    ("Superadmin mengawasi", PURPLE),
    ("Langganan & pembayaran", BRAND),
    ("Tanpa integrasi polisi", DANGER),
]
bx = 0.85
for label, col in bl:
    w = 0.30 + len(label) * 0.085
    pill(s, bx, 5.22, w, 0.32, label, fill=SURFACE2, color=col, size=9.5)
    bx += w + 0.18

footnote(s, "Setiap anggota baru — termasuk yang memakai kode undangan — wajib disetujui admin sebelum bisa memakai aplikasi.")

# ================================================================ SLIDE 3
s = slide()
header(s, "02  ·  ONBOARDING", "Cara Gabung atau Buat Lingkungan", accent=INFO)

box(s, 0.58, 1.42, 2.30, 0.80, "Daftar", "Pilih bahasa lebih dulu",
    fill=SURFACE, border=LINE)
arrow(s, 2.96, 1.68, 0.36, color=TEXT3)

box(s, 3.42, 1.42, 3.10, 0.80, "Pilih cara bergabung",
    "3 jalur di bawah ini", fill=SURFACE2, border=INFO, title_color=INFO)

paths = [
    (0.58, "JALUR 1", "Kode Undangan / QR", BRAND,
     "Ketik kode 6 huruf\natau pindai QR",
     "Kode diperiksa\n(masa berlaku, kuota)",
     "Lingkungan terisi\notomatis + peran\nyang diusulkan"),
    (4.55, "JALUR 2", "Cari Lingkungan", INFO,
     "Cari nama / kota /\nalamat",
     "Pilih dari hasil\npencarian",
     "Tulis pesan\nuntuk admin"),
    (8.52, "JALUR 3", "Buat Lingkungan Baru", PURPLE,
     "Isi nama, alamat,\nkota",
     "Lingkungan dibuat\n+ percobaan 14 hari",
     "Langsung AKTIF\nsebagai ADMIN"),
]

for x, kicker, title, col, a, b, c in paths:
    pill(s, x, 2.52, 1.05, 0.28, kicker, fill=col, color=BG, size=9)
    box(s, x, 2.90, 3.32, 0.62, title, fill=SURFACE2, border=col,
        title_color=col, title_size=13)
    box(s, x, 3.62, 3.32, 0.62, "", a, fill=SURFACE, border=LINE,
        title_size=1, body_size=10)
    box(s, x, 4.34, 3.32, 0.62, "", b, fill=SURFACE, border=LINE,
        title_size=1, body_size=10)
    box(s, x, 5.06, 3.32, 0.68, "", c, fill=SURFACE, border=LINE,
        title_size=1, body_size=10)
    for yy in (3.52, 4.24, 4.96):
        down_arrow(s, x + 1.56, yy, 0.12, color=TEXT3, w=0.16)

# outcome
box(s, 0.58, 5.94, 7.29, 0.56, "Masuk antrean → menunggu persetujuan admin",
    fill=RGBColor(0x2A, 0x22, 0x14), border=WARN, title_color=WARN, title_size=12)
box(s, 8.52, 5.94, 3.32, 0.56, "Langsung aktif (pendiri)",
    fill=RGBColor(0x14, 0x2A, 0x22), border=BRAND, title_color=BRAND, title_size=12)

footnote(s, "Kode undangan hanya MENGUSULKAN peran dan mengisi lingkungan — tidak melewati antrean persetujuan.", color=WARN)

# ================================================================ SLIDE 4
s = slide()
header(s, "03  ·  VERIFIKASI", "Persetujuan Admin & Penetapan Peran", accent=PURPLE)

box(s, 0.58, 1.55, 2.45, 1.30, "Pendaftar baru",
    "Nama · HP · Email\nAlamat rumah\nCara bergabung\nPesan ke admin",
    fill=SURFACE, border=LINE)
arrow(s, 3.11, 2.10, 0.40, color=TEXT3)

box(s, 3.59, 1.55, 2.75, 1.30, "Admin meninjau",
    "Melihat: lewat kode\nundangan / pencarian,\nkode dipakai, pesan,\ndan data pendaftar",
    fill=SURFACE2, border=PURPLE, title_color=PURPLE)

# decision
dia = s.shapes.add_shape(MSO_SHAPE.DIAMOND, inches(6.68), inches(1.50),
                         inches(1.62), inches(1.40))
dia.fill.solid()
dia.fill.fore_color.rgb = SURFACE2
dia.line.color.rgb = WARN
dia.line.width = Pt(1.5)
dia.shadow.inherit = False
tf = dia.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Setuju?"
r.font.name = FONT
r.font.size = Pt(12)
r.font.bold = True
r.font.color.rgb = WARN

arrow(s, 6.42, 2.10, 0.22, color=TEXT3)

# accept branch
textbox(s, 8.36, 2.16, 0.5, 0.22, "YA", size=9.5, color=BRAND, bold=True)
arrow(s, 8.36, 1.92, 0.46, color=BRAND)
box(s, 8.88, 1.30, 3.54, 0.52, "Terima — pilih peran",
    fill=SURFACE, border=BRAND, title_color=BRAND, title_size=12)

roles = [
    (8.88, "WARGA", BRAND, "Panik, lapor,\npeta, tamu"),
    (10.06, "SATPAM", INFO, "+ patroli,\nbuku tamu"),
    (11.24, "ADMIN", PURPLE, "+ setujui,\narea, siaran"),
]
for x, label, col, desc in roles:
    pill(s, x, 1.94, 1.10, 0.28, label, fill=col, color=BG, size=9)
    box(s, x, 2.30, 1.10, 0.66, "", desc, fill=SURFACE, border=LINE,
        title_size=1, body_size=8.5)

# reject branch
textbox(s, 7.70, 3.02, 1.0, 0.24, "TIDAK", size=9.5, color=DANGER, bold=True)
down_arrow(s, 7.42, 2.96, 0.40, color=DANGER)
box(s, 6.30, 3.44, 2.42, 0.62, "Tolak + alasan",
    "Pendaftar diberi tahu", fill=SURFACE, border=DANGER, title_color=DANGER,
    title_size=12, body_size=9.5)

# result row
box(s, 8.88, 3.16, 3.54, 0.78, "Anggota AKTIF",
    "Tombol daftar hilang → diganti sapaan pribadi",
    fill=RGBColor(0x14, 0x2A, 0x22), border=BRAND, title_color=BRAND,
    title_size=12, body_size=9.5)

# admin powers
box(s, 0.58, 4.30, 5.62, 2.02, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 0.85, 4.48, 5.0, 0.28, "WEWENANG ADMIN", size=10, color=PURPLE, bold=True)
powers = [
    "Terima / tolak setiap pendaftar",
    "Tetapkan peran: Warga · Satpam · Admin",
    "Ubah peran & tangguhkan anggota kapan saja",
    "Buat kode + QR undangan (kuota, masa berlaku, cabut)",
    "Gambar area lingkungan di peta untuk semua anggota",
    "Kirim siaran darurat + minta konfirmasi keselamatan",
]
py = 4.80
for t in powers:
    textbox(s, 0.98, py, 5.1, 0.24, "▸  " + t, size=10, color=TEXT2)
    py += 0.245

box(s, 6.44, 4.30, 5.98, 2.02, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 6.72, 4.48, 5.4, 0.28, "PENGAWASAN SUPERADMIN", size=10, color=WARN, bold=True)
sup = [
    "Akun tetap: tarafk1972@gmail.com",
    "Memantau seluruh lingkungan & admin",
    "Verifikasi pembayaran langganan",
    "Menjawab tiket bantuan (customer service)",
    "Perpanjang masa percobaan / tangguhkan lingkungan",
    "Melihat catatan aktivitas (audit log)",
]
py = 4.80
for t in sup:
    textbox(s, 6.85, py, 5.4, 0.24, "▸  " + t, size=10, color=TEXT2)
    py += 0.245

footnote(s, "Peran yang diusulkan kode undangan otomatis terpilih, tetapi admin tetap bebas mengubahnya sebelum menerima.")

# ================================================================ SLIDE 5
s = slide()
header(s, "04  ·  DARURAT (MVP)", "Alur Tombol Panik", accent=DANGER)

# big button
btn = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(0.62), inches(1.72),
                         inches(1.92), inches(1.92))
btn.fill.solid()
btn.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)
btn.line.color.rgb = DANGER
btn.line.width = Pt(2.5)
btn.shadow.inherit = False
tf = btn.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "DARURAT"
r.font.name = FONT
r.font.size = Pt(17)
r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "tahan 2 detik"
r2.font.name = FONT
r2.font.size = Pt(9)
r2.font.color.rgb = RGBColor(0xFF, 0xD0, 0xD6)

textbox(s, 0.62, 3.78, 1.92, 0.5, "Satu layar,\nsatu tombol", size=10,
        color=TEXT3, align=PP_ALIGN.CENTER, spacing=1.2)

arrow(s, 2.66, 2.56, 0.38, color=DANGER)

# capture
box(s, 3.10, 1.62, 3.05, 2.30, "", fill=SURFACE, border=DANGER, radius=0.05)
textbox(s, 3.28, 1.78, 2.7, 0.26, "YANG DIKIRIM", size=10, color=DANGER, bold=True)
items = [
    "Lokasi GPS + akurasi",
    "Lokasi langsung (terus)",
    "Profil & data medis",
    "Rekaman suara 15 detik",
    "Foto / video (opsional)",
    "Jenis darurat (opsional)",
    "Waktu kejadian",
]
iy = 2.10
for t in items:
    textbox(s, 3.30, iy, 2.75, 0.24, "•  " + t, size=10, color=TEXT2)
    iy += 0.245

arrow(s, 6.27, 2.56, 0.38, color=DANGER)

# recipients
box(s, 6.73, 1.62, 2.55, 2.30, "", fill=SURFACE, border=BRAND, radius=0.05)
textbox(s, 6.92, 1.78, 2.3, 0.26, "DIKIRIM KE", size=10, color=BRAND, bold=True)
rec = [
    ("Keluarga", BRAND),
    ("Teman terpercaya", INFO),
    ("Responder komunitas", PURPLE),
    ("Satpam", WARN),
    ("Relawan", RGBColor(0x5E, 0xEA, 0xD4)),
]
ry = 2.16
for label, col in rec:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(6.94), inches(ry + 0.055),
                           inches(0.12), inches(0.12))
    d.fill.solid()
    d.fill.fore_color.rgb = col
    d.line.fill.background()
    d.shadow.inherit = False
    textbox(s, 7.14, ry, 2.1, 0.24, label, size=10.5, color=TEXT2)
    ry += 0.31

arrow(s, 9.40, 2.56, 0.38, color=BRAND)

# response
box(s, 9.86, 1.62, 2.56, 2.30, "", fill=SURFACE, border=INFO, radius=0.05)
textbox(s, 10.04, 1.78, 2.3, 0.26, "RESPONS", size=10, color=INFO, bold=True)
resp = [
    "Penerima lihat lokasi",
    "Tekan “Saya menuju lokasi”",
    "Percakapan dua arah",
    "Pengirim lihat siapa datang",
    "Selesai: “Saya sudah aman”",
]
ry = 2.16
for t in resp:
    textbox(s, 10.06, ry, 2.3, 0.28, "•  " + t, size=10, color=TEXT2, spacing=1.0)
    ry += 0.325

# safety rail
box(s, 0.62, 4.20, 5.78, 1.28, "", fill=RGBColor(0x2A, 0x14, 0x18),
    border=DANGER, radius=0.05)
textbox(s, 0.88, 4.38, 5.2, 0.26, "PENGAMAN ALARM PALSU", size=10,
        color=DANGER, bold=True)
guards = [
    "Wajib ditahan 2 detik — sentuhan tak sengaja tidak mengirim",
    "Lepas jari sebelum penuh → batal total",
    "Setelah terkirim: “Alarm palsu — batalkan” & “Saya sudah aman”",
]
gy = 4.70
for t in guards:
    textbox(s, 1.00, gy, 5.3, 0.24, "▸  " + t, size=10, color=TEXT2)
    gy += 0.25

box(s, 6.64, 4.20, 5.78, 1.28, "", fill=RGBColor(0x14, 0x1E, 0x2C),
    border=INFO, radius=0.05)
textbox(s, 6.90, 4.38, 5.2, 0.26, "URUTAN YANG DISENGAJA", size=10,
        color=INFO, bold=True)
order = [
    "GPS diambil DULU, sebelum peringatan dikirim",
    "Rekaman suara mulai SETELAH peringatan keluar",
    "Jenis darurat dipilih SETELAH terkirim — tidak memperlambat",
]
gy = 4.70
for t in order:
    textbox(s, 7.02, gy, 5.3, 0.24, "▸  " + t, size=10, color=TEXT2)
    gy += 0.25

box(s, 0.62, 5.66, 11.80, 0.62, "TIDAK ADA INTEGRASI POLISI / 110 / 112",
    "Peringatan hanya mengalir ke jaringan warga. Pengguna diingatkan untuk menghubungi pihak berwenang sendiri bila diperlukan.",
    fill=SURFACE2, border=WARN, title_color=WARN, title_size=11.5, body_size=9.5)

footnote(s, "Mikrofon ditolak atau GPS gagal tidak pernah membatalkan pengiriman peringatan.", y=6.46)

# ================================================================ SLIDE 6
s = slide()
header(s, "05  ·  JARINGAN BANTUAN", "Siapa Menerima & Bagaimana Merespons", accent=BRAND)

# left: who
box(s, 0.58, 1.48, 5.40, 2.56, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 0.85, 1.66, 5.0, 0.26, "SIAPA YANG MENERIMA PERINGATAN", size=10,
        color=BRAND, bold=True)

groups = [
    ("PRIBADI", "Keluarga · Teman terpercaya",
     "Ditambahkan sendiri oleh tiap anggota.\nHanya menerima peringatan pemiliknya.", BRAND),
    ("KOMUNITAS", "Responder · Satpam · Relawan",
     "Wajib diverifikasi admin dulu.\nSatpam & admin otomatis termasuk.", PURPLE),
]
gy = 2.02
for kicker, title, desc, col in groups:
    pill(s, 0.85, gy, 1.28, 0.26, kicker, fill=col, color=BG, size=8.5)
    textbox(s, 2.25, gy + 0.01, 3.6, 0.24, title, size=11, color=TEXT, bold=True)
    textbox(s, 0.87, gy + 0.32, 5.0, 0.42, desc, size=9.5, color=TEXT2, spacing=1.12)
    gy += 0.86

textbox(s, 0.87, 3.70, 5.0, 0.24,
        "Aturan: kontak pribadi satu anggota tidak pernah bocor ke anggota lain.",
        size=9.5, color=TEXT3, italic=True)

# right: flow
box(s, 6.24, 1.48, 6.18, 2.56, "", fill=SURFACE, border=LINE, radius=0.05)
textbox(s, 6.52, 1.66, 5.6, 0.26, "ALUR RESPONS", size=10, color=INFO, bold=True)

steps = [
    ("1", "Peringatan tiba", "Muncul di layar utama penerima"),
    ("2", "“Saya menuju lokasi”", "Status berubah, pengirim melihatnya"),
    ("3", "Koordinasi", "Percakapan + foto di dalam insiden"),
    ("4", "Ditutup", "Pengirim tekan “Saya sudah aman”"),
]
sy = 2.04
for num, title, desc in steps:
    pill(s, 6.52, sy + 0.02, 0.28, 0.28, num, fill=SURFACE2, color=INFO, size=9.5)
    textbox(s, 6.94, sy, 5.2, 0.24, title, size=11, color=TEXT, bold=True)
    textbox(s, 6.94, sy + 0.235, 5.2, 0.24, desc, size=9.5, color=TEXT2)
    sy += 0.50

# broadcast row
box(s, 0.58, 4.24, 11.84, 0.36, "SITUASI BESAR: ADMIN KIRIM SIARAN KE SEMUA ANGGOTA",
    fill=SURFACE2, border=LINE, title_color=TEXT2, title_size=10)

bcast = [
    (0.58, "Admin kirim siaran", "Pilih tingkat:\nInformasi / Peringatan / Darurat", WARN),
    (3.54, "Instruksi keselamatan", "Contoh: “Tetap di dalam\nrumah, kunci pintu.”", INFO),
    (6.50, "Warga menjawab", "“Saya aman” atau\n“Butuh bantuan”", BRAND),
    (9.46, "Admin lihat rekap", "Siapa aman, siapa butuh\nbantuan, siapa belum jawab", PURPLE),
]
for x, title, desc, col in bcast:
    box(s, x, 4.76, 2.62, 1.16, title, desc, fill=SURFACE, border=col,
        title_color=col, title_size=11.5, body_size=9.5)
for x in (3.54, 6.50, 9.46):
    arrow(s, x - 0.30, 5.25, 0.26, color=TEXT3)

footnote(s, "Siaran dipakai untuk kejadian yang menyangkut banyak orang; tombol panik untuk darurat pribadi.", y=6.14)

# ================================================================ SLIDE 7
s = slide()
header(s, "06  ·  RINGKASAN", "Peran & Hak Akses", accent=WARN)

cols = [
    (0.58, "WARGA", BRAND, [
        "Tombol panik darurat",
        "Kirim laporan kejadian",
        "Kirim info anonim",
        "Lihat peta & area",
        "Kelola jaringan bantuan",
        "Profil darurat medis",
        "Terima siaran & konfirmasi aman",
    ]),
    (3.54, "SATPAM", INFO, [
        "Semua hak Warga",
        "Buku tamu masuk/keluar",
        "Patroli + titik pantau",
        "Tangani laporan warga",
        "Terima semua peringatan",
        "Lihat profil darurat pelapor",
        "“Saya menuju lokasi”",
    ]),
    (6.50, "ADMIN", PURPLE, [
        "Semua hak Satpam",
        "Terima / tolak pendaftar",
        "Tetapkan & ubah peran",
        "Buat kode + QR undangan",
        "Gambar area di peta",
        "Kirim siaran darurat",
        "Kelola langganan",
    ]),
    (9.46, "SUPERADMIN", WARN, [
        "tarafk1972@gmail.com",
        "Pantau semua lingkungan",
        "Awasi seluruh admin",
        "Verifikasi pembayaran",
        "Customer service (tiket)",
        "Perpanjang percobaan",
        "Catatan aktivitas",
    ]),
]

for x, title, col, items in cols:
    box(s, x, 1.48, 2.96, 0.52, title, fill=SURFACE2, border=col,
        title_color=col, title_size=14)
    box(s, x, 2.08, 2.96, 3.30, "", fill=SURFACE, border=LINE, radius=0.05)
    iy = 2.26
    for t in items:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(x + 0.17),
                               inches(iy + 0.075), inches(0.085), inches(0.085))
        d.fill.solid()
        d.fill.fore_color.rgb = col
        d.line.fill.background()
        d.shadow.inherit = False
        textbox(s, x + 0.34, iy, 2.5, 0.42, t, size=9.5, color=TEXT2, spacing=1.1)
        iy += 0.445

box(s, 0.58, 5.56, 11.84, 0.86, "",
    fill=SURFACE2, border=LINE, radius=0.05)
textbox(s, 0.85, 5.70, 11.3, 0.24, "PRINSIP UTAMA", size=10, color=WARN, bold=True)
textbox(s, 0.85, 5.94, 11.3, 0.44,
        "Warga pertama otomatis jadi Admin  ·  Semua pendaftar wajib disetujui admin\n"
        "Percobaan gratis 14 hari  ·  Peringatan hanya ke jaringan warga, bukan polisi",
        size=10.5, color=TEXT2, spacing=1.15)

footnote(s, "Warga Jaga Warga — Versi 1 (MVP)", y=6.58)

# ---------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Alur-Kerja-WJW.pptx")
prs.save(out)
print("saved:", out, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
