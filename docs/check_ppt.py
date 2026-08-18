"""
Sanity-checks Alur-Kerja-WJW.pptx without a renderer:
  - every shape sits inside the slide
  - text boxes do not overlap each other
  - rough text-overflow estimate per shape
"""

import os
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
HERE = os.path.dirname(os.path.abspath(__file__))
PATH = os.path.join(HERE, "Alur-Kerja-WJW.pptx")

prs = Presentation(PATH)
SW = prs.slide_width / EMU_IN
SH = prs.slide_height / EMU_IN

problems = []


def rect(sp):
    return (
        sp.left / EMU_IN,
        sp.top / EMU_IN,
        (sp.left + sp.width) / EMU_IN,
        (sp.top + sp.height) / EMU_IN,
    )


def overlap(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix = min(ax2, bx2) - max(ax1, bx1)
    iy = min(ay2, by2) - max(ay1, by1)
    if ix <= 0.02 or iy <= 0.02:
        return 0.0
    return ix * iy


def text_of(sp):
    if not sp.has_text_frame:
        return ""
    return "\n".join(p.text for p in sp.text_frame.paragraphs).strip()


def max_pt(sp):
    m = 0
    for p in sp.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                m = max(m, r.font.size.pt)
    return m or 12


for idx, slide in enumerate(prs.slides, start=1):
    shapes = list(slide.shapes)

    # 1. bounds
    for sp in shapes:
        if sp.left is None:
            continue
        # decorative shapes may intentionally bleed off the slide edge
        if (sp.name or "").startswith("decor-"):
            continue
        x1, y1, x2, y2 = rect(sp)
        if x1 < -0.05 or y1 < -0.05 or x2 > SW + 0.05 or y2 > SH + 0.05:
            problems.append(
                f"slide {idx}: '{text_of(sp)[:34]}' out of bounds "
                f"({x1:.2f},{y1:.2f})-({x2:.2f},{y2:.2f})"
            )

    # 2. text overflow (rough): chars that fit vs chars present
    for sp in shapes:
        t = text_of(sp)
        if not t or sp.width is None:
            continue
        pt = max_pt(sp)
        w_in = sp.width / EMU_IN
        h_in = sp.height / EMU_IN
        # ~0.50 * font-size average glyph width for Segoe UI
        char_w = (pt * 0.50) / 72
        line_h = (pt * 1.32) / 72
        per_line = max(1, int(w_in / char_w))
        lines_avail = max(1, int(h_in / line_h))
        need = 0
        for para in t.split("\n"):
            need += max(1, -(-len(para) // per_line))
        if need > lines_avail:
            problems.append(
                f"slide {idx}: text may overflow in '{t[:34]}' "
                f"(needs ~{need} lines, fits ~{lines_avail})"
            )

    # 3. overlapping *text-bearing* shapes (ignore backgrounds/containers)
    texted = [
        sp
        for sp in shapes
        if text_of(sp) and sp.width and (sp.width / EMU_IN) * (sp.height / EMU_IN) < 12
    ]
    for i in range(len(texted)):
        for j in range(i + 1, len(texted)):
            a, b = texted[i], texted[j]
            ov = overlap(rect(a), rect(b))
            if ov <= 0:
                continue
            area_a = (a.width / EMU_IN) * (a.height / EMU_IN)
            area_b = (b.width / EMU_IN) * (b.height / EMU_IN)
            frac = ov / min(area_a, area_b)
            if frac > 0.30:
                problems.append(
                    f"slide {idx}: '{text_of(a)[:24]}' overlaps "
                    f"'{text_of(b)[:24]}' ({frac:.0%})"
                )

print(f"{PATH}")
print(f"{len(prs.slides._sldIdLst)} slides · {SW:.2f}x{SH:.2f} in")
if problems:
    print(f"\n{len(problems)} issue(s):")
    for p in problems:
        print("  -", p)
    sys.exit(1)
print("\nOK: no bounds, overflow or overlap issues detected.")
